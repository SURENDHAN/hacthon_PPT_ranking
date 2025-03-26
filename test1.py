import streamlit as st
from langchain_google_genai import ChatGoogleGenerativeAI
# Immediately set page config before any other Streamlit commands
st.set_page_config(page_title="Enhanced PPT Evaluator", layout="wide")

# Correct import from langchain_community instead of langchain
from langchain_community.chat_models import ChatOllama
import matplotlib.pyplot as plt
import numpy as np
import re
from pptx import Presentation
import json
from io import BytesIO
import time
import pandas as pd
import concurrent.futures
from functools import partial

# Add custom CSS
st.markdown("""
<style>
.main .block-container {
    padding-top: 2rem;
    padding-bottom: 2rem;
}
h1, h2, h3 {
    color: #1e3a8a;
}
.stProgress > div > div {
    background-color: #06d6a0;
}
.score-box {
    text-align: center;
    margin-bottom: 15px;
}
.score-label {
    font-weight: bold;
    font-size: 1.2em;
}
.score-value {
    margin: 10px 0;
    padding: 10px;
    border-radius: 5px;
    font-size: 1.5em;
    font-weight: bold;
    color: white;
}
.weight-slider {
    padding: 10px 0;
}
.weighted-score {
    font-size: 56px;
    font-weight: bold;
    text-align: center;
    color: #0284c7;
}
.weight-container {
    background-color: #f8fafc;
    padding: 15px;
    border-radius: 10px;
    margin-bottom: 20px;
}
.rank-1 {
    background-color: gold;
    color: black;
    font-weight: bold;
}
.rank-2 {
    background-color: silver;
    color: black;
    font-weight: bold;
}
.rank-3 {
    background-color: #cd7f32;
    color: white;
    font-weight: bold;
}
.ranking-table th {
    background-color: #1e3a8a;
    color: white;
}
.file-info {
    font-size: 0.9em;
    color: #4b5563;
}
</style>
""", unsafe_allow_html=True)


@st.cache_resource(ttl=3600) 
def load_llm_model():
    try:
        return ChatGoogleGenerativeAI(
            model="gemini-1.5-pro",
            temperature=0.1,
            google_api_key='AIzaSyDlYC_kPXmiDCd6Xj9EvdAzxE7X0R' 
        )
    except Exception as e:
        st.error(f"Failed to load LLM model: {e}")
        return None

chat_model = load_llm_model()

# Optimized function to extract text from a PowerPoint file
def extract_text_from_pptx(file):
    try:
        ppt = Presentation(file)
        text = []
        slide_count = len(ppt.slides)
        
        # Create a progress bar for extraction
        progress_bar = st.progress(0)
        
        # Optimize by processing slides in chunks
        chunk_size = max(1, slide_count // 10)  # Process 10% at a time
        
        for chunk_start in range(0, slide_count, chunk_size):
            chunk_end = min(chunk_start + chunk_size, slide_count)
            chunk_slides = []
            
            for i in range(chunk_start, chunk_end):
                slide = ppt.slides[i]
                slide_text = []
                
                # More efficient text extraction by type checking first
                for shape in slide.shapes:
                    if not hasattr(shape, "text"):
                        continue
                        
                    text_frame = shape.text.strip()
                    if text_frame:
                        slide_text.append(text_frame)
                
                if slide_text:
                    chunk_slides.append(f"[SLIDE {i+1}/{slide_count}]\n" + "\n".join(slide_text))
            
            # Add chunk to the complete text
            text.extend(chunk_slides)
            
            # Update progress
            progress_bar.progress(chunk_end / slide_count)
        
        progress_bar.empty()
        
        # Join only at the end
        return "\n\n".join(text)
    except Exception as e:
        st.error(f"Error extracting text from PowerPoint: {e}")
        return ""

# Optimized extraction of scores
def extract_scores(response_text):
    """Optimized score extraction with better performance"""
    try:
        # First attempt: Try to extract JSON if the model returned it
        json_match = re.search(r'```json\s*(.*?)\s*```', response_text, re.DOTALL)
        
        if json_match:
            try:
                result = json.loads(json_match.group(1))
                return result.get('scores', {}), result.get('justifications', {})
            except json.JSONDecodeError:
                pass  # Fall back to regex if JSON parsing fails
        
        # Initialize containers
        categories = [
            "Innovation", "Technical Feasibility", "Impact", 
            "Technical Depth", "Presentation Quality", "Scalability", "Market Potential"
        ]
        scores = {}
        justifications = {}
        
        # Compile regex patterns once (more efficient)
        score_patterns = [
            re.compile(rf"{category}.*?\(Score\s*(\d+(?:\.\d+)?)/10\)", re.IGNORECASE | re.DOTALL),
            re.compile(rf"{category}.*?:\s*(\d+(?:\.\d+)?)/10", re.IGNORECASE | re.DOTALL),
            re.compile(rf"{category}.*?(\d+(?:\.\d+)?)\s*/\s*10", re.IGNORECASE | re.DOTALL)
        ]
        
        justification_patterns = [
            re.compile(rf"{category}.*?\(Score\s*\d+(?:\.\d+)?/10\)\s*(.*?)(?=\n\n|\n\*|\Z)", re.IGNORECASE | re.DOTALL),
            re.compile(rf"{category}.*?:\s*(?:\d+(?:\.\d+)?/10)\s*(.*?)(?=\n\n|\n\*|\Z)", re.IGNORECASE | re.DOTALL),
            re.compile(rf"Justification for {category}:\s*(.*?)(?=\n\n|\n\*|\Z)", re.IGNORECASE | re.DOTALL)
        ]
        
        # Process each category
        for category in categories:
            # Extract score
            for i, pattern in enumerate(score_patterns):
                score_match = pattern.search(response_text)
                if score_match:
                    try:
                        scores[category] = float(score_match.group(1))
                        break
                    except (ValueError, IndexError):
                        continue
            
            # Extract justification
            for i, pattern in enumerate(justification_patterns):
                just_match = pattern.search(response_text)
                if just_match:
                    justifications[category] = just_match.group(1).strip()
                    break
            
            # If no justification found, set a default
            if category not in justifications:
                justifications[category] = "No justification found."
        
        return scores, justifications
    except Exception as e:
        st.error(f"⚠ Error in extracting evaluation results: {e}")
        return {}, {}

# Calculate weighted average score
def calculate_weighted_average(scores, weights):
    valid_scores = {k: v for k, v in scores.items() if v is not None}
    if not valid_scores:
        return 0
    
    total_weighted_score = sum(valid_scores[category] * weights[category] for category in valid_scores)
    total_weight = sum(weights[category] for category in valid_scores)
    
    return total_weighted_score / total_weight if total_weight > 0 else 0

# Calculate the simple average score (unweighted)
def calculate_average(scores):
    valid_scores = [v for v in scores.values() if v is not None]
    return sum(valid_scores) / len(valid_scores) if valid_scores else 0

# Enhanced visualization function with weight controls
def plot_evaluation(scores, justifications, weights=None):
    valid_scores = {k: v for k, v in scores.items() if v is not None}
    
    if not valid_scores:
        st.error("❌ No valid scores available for visualization.")
        return
    
    # Initialize weights if not provided
    if weights is None:
        weights = {category: 1.0 for category in valid_scores}
    
    # Calculate weighted score
    weighted_avg = calculate_weighted_average(valid_scores, weights)
    regular_avg = calculate_average(valid_scores)
    
    # Display both scores
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown(f"""
        <div style='background-color:#f0f9ff; padding:15px; border-radius:10px; 
             text-align:center; margin:20px 0;'>
            <h3 style='margin:0; color:#0c4a6e;'>Weighted Score</h3>
            <div class='weighted-score'>
                {weighted_avg:.1f}<span style='font-size:28px;'>/10</span>
            </div>
        </div>
        """, unsafe_allow_html=True)
    
    with col2:
        st.markdown(f"""
        <div style='background-color:#f0f2f5; padding:15px; border-radius:10px; 
             text-align:center; margin:20px 0;'>
            <h3 style='margin:0; color:#4b5563;'>Unweighted Score</h3>
            <div class='weighted-score' style='color:#4b5563;'>
                {regular_avg:.1f}<span style='font-size:28px;'>/10</span>
            </div>
        </div>
        """, unsafe_allow_html=True)
    
    # Create tabs for different visualizations
    chart_tab1, chart_tab2 = st.tabs(["Bar Chart", "Radar Chart"])
    
    # Tab 1: Enhanced Bar Chart
    with chart_tab1:
        fig, ax = plt.subplots(figsize=(10, 6))
        categories = list(valid_scores.keys())
        values = list(valid_scores.values())
        
        # Color coding based on score ranges
        colors = ['#ff6b6b' if v < 4 else '#ffd166' if v < 7 else '#06d6a0' for v in values]
        
        # Create horizontal bar chart
        bars = ax.barh(categories, values, color=colors)
        
        # Add value labels
        for bar in bars:
            width = bar.get_width()
            ax.text(width + 0.1, bar.get_y() + bar.get_height()/2, f'{width:.1f}', 
                    ha='left', va='center', fontweight='bold')
        
        # Customize chart
        ax.set_xlabel("Scores (out of 10)", fontsize=12)
        ax.set_xlim(0, 10.5)
        ax.set_title("Presentation Evaluation Scores", fontsize=14, fontweight='bold')
        ax.grid(axis='x', linestyle='--', alpha=0.7)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        
        # Add vertical lines for both average scores
        ax.axvline(x=regular_avg, color='#4a4e69', linestyle='--', linewidth=2)
        ax.text(regular_avg, -0.5, f'Avg: {regular_avg:.1f}', ha='center', va='top', 
                fontweight='bold', color='#4a4e69')
                
        ax.axvline(x=weighted_avg, color='#0284c7', linestyle='-', linewidth=2)
        ax.text(weighted_avg, -0.8, f'Weighted: {weighted_avg:.1f}', ha='center', va='top', 
                fontweight='bold', color='#0284c7')
        
        st.pyplot(fig)
    
    # Tab 2: Radar Chart
    with chart_tab2:
        fig = plt.figure(figsize=(8, 8))
        ax = fig.add_subplot(111, polar=True)
        
        # Prepare data for radar chart
        categories = list(valid_scores.keys())
        values = list(valid_scores.values())
        
        # Number of categories
        N = len(categories)
        
        # Compute angle for each category
        angles = [n / float(N) * 2 * np.pi for n in range(N)]
        
        # Close the loop
        values += values[:1]
        angles += angles[:1]
        categories += categories[:1]
        
        # Plot data
        ax.plot(angles, values, linewidth=2, linestyle='solid', color='#3498db')
        
        # Fill area
        ax.fill(angles, values, alpha=0.25, color='#3498db')
        
        # Set category labels
        plt.xticks(angles[:-1], categories[:-1], fontsize=12)
        
        # Set radial limits
        ax.set_ylim(0, 10)
        plt.yticks(np.arange(2, 11, 2), ['2', '4', '6', '8', '10'], fontsize=10)
        
        # Add title
        plt.title('Evaluation Radar Chart', size=14, fontweight='bold', y=1.1)
        
        st.pyplot(fig)
    
    # Display scores and justifications as regular sections
    st.markdown("### 📋 Detailed Evaluation")
    
    for category in scores:
        col1, col2, col3 = st.columns([1, 3, 1])
        with col1:
            # Create a gauge-like visualization for the score
            score_val = scores[category] if scores[category] is not None else 0
            score_color = '#ff6b6b' if score_val < 4 else '#ffd166' if score_val < 7 else '#06d6a0'
            st.markdown(f"""
            <div class="score-box">
                <div class="score-label">{category}</div>
                <div class="score-value" style="background-color: {score_color};">
                    {score_val}/10
                </div>
                <div style="font-size: 0.9em; margin-top: 5px;">
                    Weight: {weights[category]}x
                </div>
            </div>
            """, unsafe_allow_html=True)
        
        with col2:
            st.markdown(f"**Justification:**")
            st.write(justifications[category])
        
        with col3:
            # Show the weighted contribution to final score
            weighted_contribution = score_val * weights[category]
            st.markdown(f"""
            <div style="text-align: center; margin-top: 25px;">
                <div style="font-size: 0.9em;">Weighted Score:</div>
                <div style="font-size: 1.5em; font-weight: bold; color: #0284c7;">
                    {score_val} × {weights[category]} = {weighted_contribution:.1f}
                </div>
            </div>
            """, unsafe_allow_html=True)
        
        st.markdown("---")

# Improved evaluation prompt
import json

def get_evaluation_prompt(ppt_text):
    """
    Generates an AI evaluation prompt for the given PPT content.
    Ensures blank or template PPTs are assigned 0 scores with proper justification.
    """

    # ✅ **Define What a Template Is**
    TEMPLATE_KEYWORDS = [
        "click to add title", "click to add text", "sample text", "lorem ipsum", "presentation template",
        "default slide", "add your content here", "company name", "date", "subtitle", "footer", "author",
        "outline", "table of contents", "background design", "placeholder"
    ]

    NON_TECHNICAL_KEYWORDS = [
        "team name", "college", "university", "team members", "team member-1", "team member-2", "team member-3", "team member-4",
        "detailed solution", "approach", "name of college", "name of university"
    ]

    TEMPLATE_DEFINITION = """
A **template PPT** is a generic presentation that contains pre-filled slides with placeholders,
sample text, or default content without any original or substantive material. These presentations 
often include:

1️⃣ **Title & Subtitle Placeholders** (e.g., "Click to add title")  
2️⃣ **Default Sample Content** (e.g., "Lorem Ipsum", "Your Text Here")  
3️⃣ **Pre-Designed Layouts Without Real Data**  
4️⃣ **Non-Technical Information** like team details, college names, etc.  
5️⃣ **Lack of Custom Implementation, Research, or Unique Problem-Solving Approach**  

A **valid PPT** should contain **technical explanations, workflows, architectures, algorithms, and real implementation details.**
"""

    # ✅ **Check If PPT Is Blank**
    if not ppt_text.strip():
        return json.dumps({
            "Error": "The given PPT is completely blank and does not contain any meaningful content.",
            "Overall Score": 0,
            "scores": {
                "Innovation": 0,
                "Technical Feasibility": 0,
                "Impact": 0,
                "Technical Depth": 0,
                "Presentation Quality": 0,
                "Scalability": 0,
                "Market Potential": 0
            },
            "justifications": "This presentation is entirely blank and does not contain any text, diagrams, or substantive content.",
            "summary": "No evaluation performed as the submitted PPT is completely blank.",
            "template_definition": TEMPLATE_DEFINITION
        }, indent=4)

    # ✅ **Check If PPT Is a Template**
    if any(keyword in ppt_text.lower() for keyword in TEMPLATE_KEYWORDS):
        return json.dumps({
            "Error": "The given PPT is a template and does not contain any original or substantive content.",
            "Overall Score": 0,
            "scores": {
                "Innovation": 0,
                "Technical Feasibility": 0,
                "Impact": 0,
                "Technical Depth": 0,
                "Presentation Quality": 0,
                "Scalability": 0,
                "Market Potential": 0
            },
            "justifications": "This presentation appears to be a generic template with sample text, placeholders, or default slides.",
            "summary": "No evaluation performed as the submitted PPT is a template without unique content.",
            "template_definition": TEMPLATE_DEFINITION
        }, indent=4)

    # ✅ **Remove Non-Technical Information Before Passing to AI**
    for keyword in NON_TECHNICAL_KEYWORDS:
        ppt_text = ppt_text.replace(keyword, "")

    # ✅ **Proceed with AI Evaluation Only If the PPT Has Content**
    return f"""
🔹 **You are an expert judge evaluating a hackathon/ideathon presentation for college students.**  
⚠ **IMPORTANT:** Focus **ONLY** on the **CORE TECHNICAL CONTENT** of the presentation.  
🚫 Ignore **introductory slides, team details, slide titles, college/university names, and non-substantive content.**  

---

## 📌 PRESENTATION CONTENT:
{ppt_text}

---

## **📝 EVALUATION TASK:**  
1️⃣ **Extract only the technical and substantive content** from the presentation.  
2️⃣ **Evaluate the presentation strictly based on the provided criteria.**  
3️⃣ **Ensure blank PPTs, templates, and weak content are handled correctly.**  
4️⃣ **FORMAT the response exactly as requested.**  

---

## **📝 EVALUATION CRITERIA (1-10 SCALE):**  

1️⃣ **Innovation**  
   - Is the idea original, creative, and unique?  
   - Does it introduce a novel approach using emerging technologies?  
   - Does it challenge traditional solutions and introduce fresh perspectives?  
   - **Penalty** if the idea already exists without meaningful improvement.  

2️⃣ **Technical Feasibility**  
   - Can the solution be implemented realistically with available tools and resources?  
   - Is the **problem statement well-defined** and does the solution fully address it?  
   - Does the architecture/workflow follow logical steps for implementation?  
   - **Penalty** if the solution lacks clear implementation details or is impractical.  

3️⃣ **Impact**  
   - Does the solution solve a **real-world problem** effectively?  
   - Are there measurable **social, economic, or environmental** benefits?  
   - Can this solution be **deployed in academic, industrial, or community settings?**  
   - **Penalty** if the impact is **minimal or unclear.**  

4️⃣ **Technical Depth**  
   - Is the approach **technically sophisticated** yet aligned with the problem?  
   - Are the chosen **technologies, algorithms, frameworks, and methodologies justified?**  
   - Does it integrate **best practices** and deep technical understanding?  
   - **Penalty** if the approach lacks **depth, justification, or methodology.**  

5️⃣ **Presentation Quality**  
   - Is the **PPT clear, well-structured, and logically organized?**  
   - Are there **technical diagrams, workflows, or architecture representations?**  
   - Are there **well-structured visuals, engaging storytelling, and concise explanations?**  
   - **Penalty** if the presentation is **confusing, cluttered, or lacks key details.**  

6️⃣ **Scalability**  
   - Can the **solution expand** beyond its initial scope?  
   - Is it **adaptable to different industries, geographies, or user bases?**  
   - Has **performance, efficiency, and resource optimization** been considered?  
   - **Penalty** if the idea is **too niche or non-scalable.**  

7️⃣ **Market Potential**  
   - Does the solution **have demand** or a target audience?  
   - Is it commercially feasible and competitive?  
   - Can it be adopted by **businesses, governments, or communities?**  
   - **Penalty** if there is **no clear market need.**  

---

## **📝 FINAL JSON OUTPUT FORMAT:**  

```json
{{
    "scores": {{
        "Innovation": X.X,
        "Technical Feasibility": X.X,
        "Impact": X.X,
        "Technical Depth": X.X,
        "Presentation Quality": X.X,
        "Scalability": X.X,
        "Market Potential": X.X
    }},
    "justifications": {{
        "Innovation": "Detailed explanation...",
        "Technical Feasibility": "Detailed explanation...",
        "Impact": "Detailed explanation...",
        "Technical Depth": "Detailed explanation...",
        "Presentation Quality": "Detailed explanation...",
        "Scalability": "Detailed explanation...",
        "Market Potential": "Detailed explanation..."
    }},
    "overall_average": X.X,
    "summary": "Brief overall assessment..."
}}
"""
# Helper function to get default weights
def get_default_weights():
    return {
    "Innovation": 4.00,
    "Technical Feasibility": 2.00,
    "Impact": 3.00,
    "Technical Depth": 1.50,
    "Presentation Quality": 2.50,
    "Scalability": 2.00,
    "Market Potential": 2.00
}

        
    

# Function to process a single presentation
def process_single_presentation(uploaded_file):
    # Create a BytesIO object from the uploaded file
    file_bytes = BytesIO(uploaded_file.getvalue())
    file_name = uploaded_file.name
    
    # Extract text
    ppt_text = extract_text_from_pptx(file_bytes)
    if not ppt_text:
        return None
        
    # Create result dictionary with extracted text
    result = {
        'file_name': file_name,
        'ppt_text': ppt_text
    }
    
    return result

# Parallel processing of PowerPoint files
def process_presentations_in_parallel(uploaded_files):
    # Create a partial function with fixed parameters
    extract_func = partial(process_single_presentation)
    
    # Create a thread pool and process files in parallel
    with concurrent.futures.ThreadPoolExecutor() as executor:
        # Submit all files for processing
        future_to_file = {executor.submit(extract_func, uploaded_file): uploaded_file 
                         for uploaded_file in uploaded_files}
        
        # Process results as they complete
        results = []
        for future in concurrent.futures.as_completed(future_to_file):
            file = future_to_file[future]
            try:
                result = future.result()
                if result:
                    results.append(result)
                    st.success(f"✅ Successfully processed {file.name}")
            except Exception as e:
                st.error(f"❌ Error processing {file.name}: {e}")
                
    return results

# Batch evaluation function
def batch_evaluate_presentations(presentation_data_list, chat_model, weights):
    """Process multiple presentations with batched AI calls"""
    results = []
    
    # Create a progress bar for the entire batch
    total = len(presentation_data_list)
    progress_text = st.empty()
    progress_bar = st.progress(0)
    
    for i, data in enumerate(presentation_data_list):
        file_name = data['file_name']
        ppt_text = data['ppt_text']
        
        # Update progress
        progress_text.text(f"Evaluating {i+1}/{total}: {file_name}")
        progress_bar.progress((i) / total)
        
        try:
            # Generate prompt
            prompt = get_evaluation_prompt(ppt_text)
            
            # Get response
            response = chat_model.invoke(prompt)
            
            if not response or not hasattr(response, "content"):
                st.error(f"❌ Received empty response from AI model for {file_name}.")
                continue
            
            # Extract scores
            response_text = response.content
            scores, justifications = extract_scores(response_text)
            
            if not scores or all(v is None for v in scores.values()):
                st.warning(f"⚠️ Could not extract proper scores from AI response for {file_name}.")
                continue
            
            # Calculate scores
            weighted_score = calculate_weighted_average(scores, weights)
            unweighted_score = calculate_average(scores)
            
            # Store results
            result = {
                'file_name': file_name,
                'scores': scores,
                'justifications': justifications,
                'weighted_score': weighted_score,
                'unweighted_score': unweighted_score,
                'response_text': response_text
            }
            
            results.append(result)
            
            # Update progress
            progress_bar.progress((i + 1) / total)
            
        except Exception as e:
            st.error(f"❌ Error during evaluation of {file_name}: {e}")
            continue
    
    # Clear progress indicators
    progress_text.empty()
    progress_bar.empty()
    
    return results

# Function to visualize rankings
def display_rankings(results_list):
    if not results_list:
        st.warning("No presentations have been evaluated yet.")
        return
    
    # Add rank property to results
    for i, result in enumerate(sorted(results_list, key=lambda x: x['weighted_score'], reverse=True)):
        result['rank'] = i + 1
    
    # Sort by rank
    sorted_results = sorted(results_list, key=lambda x: x['rank'])
    
    # Create DataFrame for display
    df = pd.DataFrame([
        {
            'Rank': result['rank'],
            'File Name': result['file_name'],
            'Weighted Score': f"{result['weighted_score']:.1f}",
            'Unweighted Score': f"{result['unweighted_score']:.1f}",
            'Innovation': f"{result['scores'].get('Innovation', 0):.1f}",
            'Technical Feasibility': f"{result['scores'].get('Technical Feasibility', 0):.1f}",
            'Impact': f"{result['scores'].get('Impact', 0):.1f}",
            'Technical Depth': f"{result['scores'].get('Technical Depth', 0):.1f}",
            'Presentation Quality': f"{result['scores'].get('Presentation Quality', 0):.1f}"
        } for result in sorted_results
    ])
    
    # Create a visualization of the rankings
    st.markdown("### 🏆 Presentation Rankings")
    
    # Use Streamlit's dataframe with styling
    st.dataframe(
        df,
        column_config={
            "Rank": st.column_config.NumberColumn(
                "🏆 Rank",
                help="Position in the ranking",
                format="%d"
            ),
            "Weighted Score": st.column_config.NumberColumn(
                "Weighted Score",
                help="Score with criteria weights applied",
                format="%.1f"
            ),
            "Unweighted Score": st.column_config.NumberColumn(
                "Raw Score",
                help="Average score without weights",
                format="%.1f"
            ),
        },
        hide_index=True,
        use_container_width=True
    )
    
    # Create a bar chart to visualize the rankings
    fig, ax = plt.subplots(figsize=(10, max(5, len(sorted_results) * 0.4)))
    
    # Get data for plotting
    names = [f"{r['file_name']}" for r in sorted_results]
    scores = [r['weighted_score'] for r in sorted_results]
    
    # Define colors based on rank
    colors = ['gold' if r['rank'] == 1 else 'silver' if r['rank'] == 2 
             else '#cd7f32' if r['rank'] == 3 else '#3498db' for r in sorted_results]
    
    # Create horizontal bar chart
    y_pos = np.arange(len(names))
    bars = ax.barh(y_pos, scores, color=colors)
    
    # Add labels
    ax.set_yticks(y_pos)
    ax.set_yticklabels(names)
    ax.invert_yaxis()  # Labels read top-to-bottom
    ax.set_xlabel('Weighted Score')
    ax.set_title('Presentation Rankings by Score')
    
    # Add value labels
    for bar in bars:
        width = bar.get_width()
        ax.text(width + 0.1, bar.get_y() + bar.get_height()/2, f'{width:.1f}', 
                ha='left', va='center', fontweight='bold')
    
    # Set x-axis limit
    ax.set_xlim(0, 10.5)
    
    # Remove spines
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    
    # Display the chart
    st.pyplot(fig)
    
    return sorted_results

# App UI and main functionality
def main():
    # App header
    st.title("📊 Multi-Presentation Evaluator & Ranker")
    st.markdown("""
    Upload multiple PowerPoint presentations to receive AI-powered evaluations 
    and see them ranked based on weighted scores.
    """)
    
    # Initialize session state
    if 'weights' not in st.session_state:
        st.session_state.weights = get_default_weights()
    
    if 'evaluated_files' not in st.session_state:
        st.session_state.evaluated_files = []
    
    # Weight configuration section
    # with st.expander("⚖️ Adjust Evaluation Criteria Weights"):
    #     st.markdown("""
    #     <div class="weight-container">
    #     <p>Adjust the importance of each criterion by moving the sliders below. 
    #     These weights will be applied to all presentations.</p>
    #     </div>
    #     """, unsafe_allow_html=True)
        
    #     # Create weight sliders in 2 columns for better layout
    #     col1, col2 = st.columns(2)
        
    #     # Reset weights button
    #     if st.button("Reset to Default Weights"):
    #         st.session_state.weights = get_default_weights()
    #         st.experimental_rerun()
        
    #     # Create sliders for each category
    #     categories = ["Innovation", "Technical Feasibility", "Impact", "Technical Depth", "Presentation Quality"]
    #     for i, category in enumerate(categories):
    #         # Alternate between columns
    #         with col1 if i % 2 == 0 else col2:
    #             st.session_state.weights[category] = st.slider(
    #                 f"{category} Weight", 
    #                 min_value=0.5, 
    #                 max_value=5.0, 
    #                 value=st.session_state.weights.get(category, 1.0),
    #                 step=0.01,
    #                 key=f"weight_{category}",
    #                 help=f"Set the importance of {category} in the overall score"
    #             )
    
    # Display current ranking if any files have been evaluated
    if st.session_state.evaluated_files:
        ranked_results = display_rankings(st.session_state.evaluated_files)
        
        # Add an option to clear all evaluations
        if st.button("Clear All Evaluations", type="primary", help="Remove all current evaluations"):
            st.session_state.evaluated_files = []
            st.rerun()

    
    # File uploader for multiple files
    uploaded_files = st.file_uploader(
        "Upload PowerPoint files (.pptx)", 
        type=["pptx"],
        accept_multiple_files=True,
        help="You can select multiple files to evaluate and rank them."
    )

    if uploaded_files:
        # Check if there are new files to evaluate
        existing_filenames = [result['file_name'] for result in st.session_state.evaluated_files]
        new_files = [f for f in uploaded_files if f.name not in existing_filenames]
        
        if new_files:
            st.markdown(f"### 🔍 Evaluating {len(new_files)} new presentations")
            
            # Check if model is available
            if not chat_model:
                st.error("❌ LLM model is not available. Please check your configuration.")
            else:
                # STEP 1: Extract text from all presentations in parallel
                with st.status("Extracting text from presentations...") as status:
                    presentation_data = process_presentations_in_parallel(new_files)
                    status.update(label=f"Text extraction complete for {len(presentation_data)} files", state="complete")
                
                # STEP 2: Batch evaluate presentations
                if presentation_data:
                    with st.status("Running AI evaluations...") as status:
                        results = batch_evaluate_presentations(presentation_data, chat_model, st.session_state.weights)
                        status.update(label=f"Evaluation complete for {len(results)} files", state="complete")
                    
                    # Add results to session state
                    st.session_state.evaluated_files.extend(results)
                    
                    # Success message
                    st.success(f"✅ Successfully evaluated {len(results)} presentations")
                    
                    # After processing all new files, display the updated rankings
                    if st.session_state.evaluated_files:
                        ranked_results = display_rankings(st.session_state.evaluated_files)
                        st.markdown("### 🎯 All Evaluations Complete")
                else:
                    st.warning("No valid presentation data was extracted. Please check your files.")
        else:
            st.info("All uploaded files have already been evaluated. View the rankings above.")
    else:
        # Display information when no file is uploaded
        if not st.session_state.evaluated_files:
            st.info("👆 Upload PowerPoint files to get started")
            
            # Show example visualization
            
            
    # Individual presentation view
    with st.expander("🔍 View Individual Presentation Details"):
        if st.session_state.evaluated_files:
            # Select a file to view details
            file_options = [result['file_name'] for result in st.session_state.evaluated_files]
            selected_file = st.selectbox("Select a presentation to view detailed evaluation:", file_options)
            
            # Find the selected file data
            selected_data = next((result for result in st.session_state.evaluated_files 
                                if result['file_name'] == selected_file), None)
            
            if selected_data:
                # Show file info and rank
                st.markdown(f"""
                <div style='background-color:#f8fafc; padding:15px; border-radius:10px; margin-bottom:20px;'>
                    <h3 style='margin:0;'>🗂️ {selected_data['file_name']}</h3>
                    <p class='file-info'>Rank: {selected_data.get('rank', 'Unranked')} | 
                    Weighted Score: {selected_data['weighted_score']:.1f}/10</p>
                </div>
                """, unsafe_allow_html=True)
                
                # Plot visualization with the current weights
                plot_evaluation(selected_data['scores'], selected_data['justifications'], 
                               st.session_state.weights)
                
                # Show raw AI response
                if st.checkbox("Show Raw AI Evaluation"):
                    st.markdown("### 🤖 Raw AI Response")
                    st.text_area("Full evaluation text:", selected_data['response_text'], 
                                height=300, key="raw_response")
        else:
            st.info("No presentations have been evaluated yet.")

if __name__ == "__main__":
    main()
