import streamlit as st
from langchain.chat_models import ChatOllama
from pptx import Presentation
import json

# Initialize Ollama model
chat_model = ChatOllama(model="gemma3:1b")

def extract_text_from_ppt(ppt_file):
    """Extracts text from a PowerPoint (.pptx) file slide by slide."""
    try:
        presentation = Presentation(ppt_file)
        text_slides = []

        for i, slide in enumerate(presentation.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                    slide_text.append(shape.text_frame.text.strip())

            if slide_text:
                text_slides.append(f"Slide {i}:\n" + "\n".join(slide_text))

        return text_slides
    except Exception as e:
        st.error(f"Error extracting text: {e}")
        return []

def get_evaluation_prompt(slide_text):
    """Creates a single prompt evaluating all criteria at once."""
    if len(slide_text) > 1500:
        slide_text = slide_text[:1500] + "..."

    return f"""
You are an expert judge evaluating a hackathon presentation slide.

Analyze the slide based on the following **criteria**:
1. **Innovation**: How novel and creative is the idea?
2. **Technical Feasibility**: How practical and achievable is it?
3. **Impact**: What is its real-world usefulness and scalability?
4. **Presentation**: Is the slide clear, well-structured, and visually appealing?
5. **Implementation**: How well is the prototype/code executed?

**Slide Content:**
{slide_text}

Provide a structured response in the format:
**Innovation**: (Score 1-10) Explanation  
**Technical Feasibility**: (Score 1-10) Explanation  
**Impact**: (Score 1-10) Explanation  
**Presentation**: (Score 1-10) Explanation  
**Implementation**: (Score 1-10) Explanation  
"""

def evaluate_slide(slide_text):
    """Evaluates a slide based on multiple criteria using a single prompt."""
    prompt = get_evaluation_prompt(slide_text)
    response = chat_model.predict(prompt)
    
    criteria = ["Innovation", "Technical Feasibility", "Impact", "Presentation", "Implementation"]
    evaluation = {}
    
    for crit in criteria:
        try:
            part = response.split(f"**{crit}**:")[1].strip()
            score, explanation = part.split(" ", 1)
            evaluation[crit] = {"Score": score.strip("()"), "Explanation": explanation.strip()}
        except Exception:
            evaluation[crit] = {"Score": "N/A", "Explanation": "Parsing error"}
    
    return evaluation

def main():
    st.title("📊 Hackathon PPT Evaluator")
    st.write("Upload a PowerPoint presentation (.pptx) and get automated slide evaluations.")

    uploaded_file = st.file_uploader("📂 Upload PPTX file", type=["pptx"])
    
    if uploaded_file is not None:
        with st.spinner("⏳ Extracting text from slides..."):
            slides = extract_text_from_ppt(uploaded_file)
        
        if slides:
            evaluations = []
            
            for i, slide_text in enumerate(slides, 1):
                with st.spinner(f"🔍 Evaluating Slide {i}..."):
                    evaluation = evaluate_slide(slide_text)
                    evaluations.append({f"Slide {i}": evaluation})
                    
                    st.subheader(f"📌 Slide {i} Evaluation")
                    
                    # Formatting the output properly
                    eval_text = ""
                    for crit, data in evaluation.items():
                        eval_text += f"**{crit}**: **{data['Score']}**\n\n> {data['Explanation']}\n\n"
                    
                    st.markdown(eval_text, unsafe_allow_html=True)
            
            # Save to JSON
            json_data = json.dumps(evaluations, indent=4)
            st.download_button("📥 Download Results as JSON", json_data, "ppt_evaluation_results.json", "application/json")
        else:
            st.error("❌ No valid slides found in the uploaded file.")

if __name__ == "__main__":
    main()
