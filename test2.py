import streamlit as st
from langchain.chat_models import ChatOllama
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from pptx import Presentation
import json
import pandas as pd

# Initialize Ollama model
chat_model = ChatOllama(model="gemma3:1b")

def extract_text_from_ppt(ppt_file):
    """Extracts text from a PowerPoint (.pptx) file slide by slide."""
    try:
        presentation = Presentation(ppt_file)
        text_slides = []
        for i, slide in enumerate(presentation.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                    slide_text.append(shape.text_frame.text.strip())
            if slide_text:
                text_slides.append(f"Slide {i}:\n" + "\n".join(slide_text))
        return text_slides
    except Exception as e:
        st.error(f"Error extracting text: {e}")
        return []

def create_chain_of_thought_prompt(slide_text):
    """Creates a chain of thought prompt for more detailed evaluation."""
    if len(slide_text) > 1500:
        slide_text = slide_text[:1500] + "..."
    
    prompt_template = """
You are an expert judge evaluating a hackathon presentation slide.
Analyze the following slide content carefully.

**Slide Content:**
{slide_text}

Please think step by step about each of these criteria:

1. INNOVATION:
   - What is the core idea presented?
   - How novel or unique is this idea compared to existing solutions?
   - Does it solve a problem in a new way?
   - Think through: What score from 1-10 would you give for innovation and why?

2. TECHNICAL FEASIBILITY:
   - What technologies are mentioned or implied?
   - How realistic is implementation given current technology?
   - Are there technical challenges that need to be addressed?
   - Think through: What score from 1-10 would you give for technical feasibility and why?

3. IMPACT:
   - Who would benefit from this solution?
   - How significant is the problem being solved?
   - What is the potential scale of impact?
   - Think through: What score from 1-10 would you give for impact and why?

4. PRESENTATION:
   - Is the slide clearly structured and easy to understand?
   - Does it effectively communicate the key points?
   - Is it visually well-organized (based on text layout)?
   - Think through: What score from 1-10 would you give for presentation and why?

5. IMPLEMENTATION:
   - Is there evidence of a working prototype or code?
   - How complete does the implementation seem?
   - Are there technical details that demonstrate depth?
   - Think through: What score from 1-10 would you give for implementation and why?

After thinking through each criterion, provide your final evaluation in this exact format:

**Innovation**: (Score 1-10) Explanation  
**Technical Feasibility**: (Score 1-10) Explanation  
**Impact**: (Score 1-10) Explanation  
**Presentation**: (Score 1-10) Explanation  
**Implementation**: (Score 1-10) Explanation  

Your final scores should reflect your careful analysis from the steps above.
"""
    
    return PromptTemplate(
        input_variables=["slide_text"],
        template=prompt_template
    )

def evaluate_slide(slide_text):
    """Evaluates a slide using chain of thought prompting."""
    prompt = create_chain_of_thought_prompt(slide_text)
    chain = LLMChain(llm=chat_model, prompt=prompt)
    response = chain.run(slide_text=slide_text)
    
    criteria = ["Innovation", "Technical Feasibility", "Impact", "Presentation", "Implementation"]
    evaluation = {}
    
    for crit in criteria:
        try:
            if f"**{crit}**:" in response:
                part = response.split(f"**{crit}**:")[1].strip()
                if "**" in part:  # If there's another criterion after this one
                    part = part.split("**")[0].strip()
                
                # Extract score and explanation
                if "(" in part and ")" in part:
                    score_part = part[part.find("(")+1:part.find(")")]
                    if "-" in score_part:
                        score = score_part.split("-")[1].strip()  # Take the latter part of "1-10"
                    else:
                        score = score_part.strip()
                    
                    explanation = part[part.find(")")+1:].strip()
                else:
                    # Fallback if format is not perfect
                    words = part.split()
                    score = words[0].strip("()")
                    explanation = " ".join(words[1:]).strip()
                
                evaluation[crit] = {"Score": score, "Explanation": explanation}
            else:
                evaluation[crit] = {"Score": "N/A", "Explanation": "Not found in response"}
        except Exception as e:
            evaluation[crit] = {"Score": "N/A", "Explanation": f"Parsing error: {str(e)}"}
    
    return evaluation, response

def main():
    st.set_page_config(page_title="Hackathon PPT Evaluator", layout="wide")
    
    st.title("📊 Hackathon PPT Evaluator")
    st.write("Upload a PowerPoint presentation (.pptx) and get detailed evaluations with chain of thought reasoning.")
    
    with st.sidebar:
        st.header("Settings")
        show_raw_response = st.checkbox("Show AI's raw responses", value=False)
    
    uploaded_file = st.file_uploader("Upload PPTX file", type=["pptx"])
    
    if uploaded_file is not None:
        with st.spinner("Extracting text from slides..."):
            slides = extract_text_from_ppt(uploaded_file)
        
        if slides:
            evaluations = []
            raw_responses = []
            
            for i, slide_text in enumerate(slides, 1):
                with st.spinner(f"Evaluating Slide {i}..."):
                    evaluation, raw_response = evaluate_slide(slide_text)
                    evaluations.append({f"Slide {i}": evaluation})
                    raw_responses.append(raw_response)
                    
                    st.subheader(f"Slide {i} Evaluation")
                    
                    # Display the content of the slide
                    with st.expander("View Slide Content"):
                        st.text(slide_text)
                    
                    # Display evaluation
                    table_data = []
                    total_score = 0
                    valid_scores = 0
                    
                    for crit, data in evaluation.items():
                        try:
                            score = float(data["Score"])
                            total_score += score
                            valid_scores += 1
                            table_data.append([crit, f"{score:.1f}/10", data["Explanation"]])
                        except (ValueError, TypeError):
                            table_data.append([crit, data["Score"], data["Explanation"]])
                    
                    # Calculate average score if possible
                    if valid_scores > 0:
                        avg_score = total_score / valid_scores
                        table_data.append(["Average", f"{avg_score:.1f}/10", ""])
                    
                    # Create a DataFrame for better display
                    df = pd.DataFrame(table_data, columns=["Criterion", "Score", "Explanation"])
                    st.dataframe(df, use_container_width=True)
                    
                    if show_raw_response:
                        with st.expander("AI's Raw Response"):
                            st.text(raw_response)
            
            # Create final summary
            if len(evaluations) > 1:
                st.subheader("Summary of All Slides")
                summary_data = []
                
                for i, eval_dict in enumerate(evaluations):
                    slide_key = f"Slide {i+1}"
                    slide_eval = eval_dict[slide_key]
                    
                    scores = []
                    for crit, data in slide_eval.items():
                        try:
                            scores.append(float(data["Score"]))
                        except (ValueError, TypeError):
                            pass
                    
                    if scores:
                        avg_score = sum(scores) / len(scores)
                        summary_data.append([slide_key, f"{avg_score:.1f}/10"])
                    else:
                        summary_data.append([slide_key, "N/A"])
                
                summary_df = pd.DataFrame(summary_data, columns=["Slide", "Average Score"])
                st.dataframe(summary_df, use_container_width=True)
            
            # Save to JSON
            json_data = json.dumps(evaluations, indent=4)
            st.download_button("Download Results as JSON", json_data, "ppt_evaluation_results.json", "application/json")
            
            # Download as CSV
            csv_data = []
            headers = ["Slide", "Criterion", "Score", "Explanation"]
            
            for i, eval_dict in enumerate(evaluations):
                slide_key = f"Slide {i+1}"
                slide_eval = eval_dict[slide_key]
                
                for crit, data in slide_eval.items():
                    csv_data.append([slide_key, crit, data["Score"], data["Explanation"]])
            
            csv_df = pd.DataFrame(csv_data, columns=headers)
            csv = csv_df.to_csv(index=False)
            st.download_button(
                "Download Results as CSV",
                csv,
                "ppt_evaluation_results.csv",
                "text/csv",
                key='download-csv'
            )
            
        else:
            st.error("No valid slides found in the uploaded file.")

if __name__ == "__main__":
    main()
